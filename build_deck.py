"""
Migration Validator — Modern Hackathon Deck Builder v2
Run: python build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Palette ────────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
DARKER     = RGBColor(0x06, 0x0D, 0x18)
CARD_BG    = RGBColor(0x16, 0x2A, 0x40)
CARD_BG2   = RGBColor(0x1E, 0x38, 0x52)
A_BLUE     = RGBColor(0x00, 0xB4, 0xFF)
A_CYAN     = RGBColor(0x00, 0xFF, 0xD4)
A_ORANGE   = RGBColor(0xFF, 0x8C, 0x00)
A_GREEN    = RGBColor(0x00, 0xE5, 0x76)
A_RED      = RGBColor(0xFF, 0x3D, 0x3D)
A_PURPLE   = RGBColor(0x9B, 0x5D, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xBD, 0xD5, 0xEE)
DIM        = RGBColor(0x6B, 0x8A, 0xA8)

W = Inches(13.33)
H = Inches(7.5)
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── XML / visual helpers ───────────────────────────────────────────────────────

def _h(c): return "{:02X}{:02X}{:02X}".format(c[0], c[1], c[2])

def apply_grad(shape, c1, c2, angle_deg=135):
    sp = shape.element.spPr
    for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'), qn('a:pattFill')]:
        for el in sp.findall(tag): sp.remove(el)
    ang = int(angle_deg * 60000) % 21600000
    xml = (f'<a:gradFill xmlns:a="{NS_A}"><a:gsLst>'
           f'<a:gs pos="0"><a:srgbClr val="{_h(c1)}"/></a:gs>'
           f'<a:gs pos="100000"><a:srgbClr val="{_h(c2)}"/></a:gs>'
           f'</a:gsLst><a:lin ang="{ang}" scaled="0"/></a:gradFill>')
    sp.insert(0, etree.fromstring(xml))

def make_round(shape, adj=8000):
    sp = shape.element.spPr
    pg = sp.find(qn('a:prstGeom'))
    if pg is None: return
    pg.set('prst', 'roundRect')
    av = pg.find(qn('a:avLst'))
    if av is None: av = etree.SubElement(pg, qn('a:avLst'))
    for gd in av.findall(qn('a:gd')): av.remove(gd)
    av.append(etree.fromstring(f'<a:gd xmlns:a="{NS_A}" name="adj" fmla="val {adj}"/>'))

def add_shadow(shape, blur=57150, dist=38100, trans=55):
    sp = shape.element.spPr
    for el in sp.findall(qn('a:effectLst')): sp.remove(el)
    alpha = int((100 - trans) * 1000)
    xml = (f'<a:effectLst xmlns:a="{NS_A}"><a:outerShdw blurRad="{blur}" dist="{dist}" '
           f'dir="2700000" algn="tl" rotWithShape="0">'
           f'<a:srgbClr val="000000"><a:alpha val="{alpha}"/></a:srgbClr>'
           f'</a:outerShdw></a:effectLst>')
    sp.append(etree.fromstring(xml))

def add_glow(shape, color, radius=38100, trans=35):
    sp = shape.element.spPr
    for el in sp.findall(qn('a:effectLst')): sp.remove(el)
    alpha = int((100 - trans) * 1000)
    xml = (f'<a:effectLst xmlns:a="{NS_A}"><a:glow rad="{radius}">'
           f'<a:srgbClr val="{_h(color)}"><a:alpha val="{alpha}"/></a:srgbClr>'
           f'</a:glow></a:effectLst>')
    sp.append(etree.fromstring(xml))

def set_alpha(shape, alpha_pct):
    sp = shape.element.spPr
    for sf in sp.findall(qn('a:solidFill')):
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha_pct * 1000)))

# ── Low-level drawing ──────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    if lw:   s.line.width = lw
    if not line: s.line.fill.background()
    return s

def oval(slide, cx, cy, r, fill=None):
    s = slide.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    return s

def rnd(slide, x, y, w, h, fill=None, grad=None, line=None, lw=None,
        adj=8000, shadow=False, glow=None):
    s = rect(slide, x, y, w, h, fill=fill or CARD_BG, line=line, lw=lw)
    make_round(s, adj)
    if grad: apply_grad(s, grad[0], grad[1], grad[2] if len(grad)>2 else 135)
    if shadow: add_shadow(s)
    if glow: add_glow(s, glow)
    return s

def txb(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
        italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tf

def addpara(tf, text, size=12, bold=False, color=LIGHT, italic=False,
            align=PP_ALIGN.LEFT, space=5):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return p

def connector(slide, x1, y1, x2, y2, color=A_BLUE, w=Pt(1.5)):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = w

# ── Composite helpers ──────────────────────────────────────────────────────────

def bg(slide, c1=DARKER, c2=DARK_BG, angle=140):
    s = rect(slide, 0, 0, W, H); apply_grad(s, c1, c2, angle); s.line.fill.background()

def orbs(slide, color=A_BLUE):
    s1 = oval(slide, W + Inches(0.5), Inches(-1), Inches(3.5), fill=color)
    set_alpha(s1, 6)
    s2 = oval(slide, Inches(-0.8), H + Inches(0.5), Inches(2.8), fill=color)
    set_alpha(s2, 5)

def grad_bar(slide, x, y, w, h, c1, c2, angle=0):
    s = rect(slide, x, y, w, h); apply_grad(s, c1, c2, angle); s.line.fill.background()

def footer(slide, n=""):
    rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=RGBColor(0x05, 0x0A, 0x12))
    txb(slide, "EPAM Hackathon 2026  ·  Stream 3: Connectors for Gemini Applications",
        Inches(0.5), H - Inches(0.33), Inches(11), Inches(0.28), size=8, color=DIM)
    if n:
        txb(slide, n, W - Inches(0.75), H - Inches(0.33), Inches(0.55), Inches(0.28),
            size=8, color=DIM, align=PP_ALIGN.RIGHT)

def pill(slide, text, x, y, bg_c, fg_c=DARKER, size=9, w=Inches(2.0), h=Inches(0.28)):
    s = rnd(slide, x, y, w, h, fill=bg_c, adj=20000)
    txb(slide, text, x, y, w, h, size=size, bold=True, color=fg_c, align=PP_ALIGN.CENTER)

def section_tag(slide, text, color, x=Inches(0.55), y=Inches(0.35)):
    pill(slide, text, x, y, bg_c=RGBColor(
        min(color[0]+20, 0x3F), min(color[1]+10, 0x1F), min(color[2]+10, 0x1F)),
        fg_c=color, size=9, w=Inches(2.2), h=Inches(0.28))

def big_title(slide, line1, line2="", color2=A_BLUE, y=Inches(0.75)):
    txb(slide, line1, Inches(0.55), y, Inches(12.3), Inches(0.85),
        size=40, bold=True, color=WHITE)
    if line2:
        txb(slide, line2, Inches(0.55), y + Inches(0.82), Inches(12.3), Inches(0.85),
            size=40, bold=True, color=color2)
    gy = y + (Inches(1.62) if line2 else Inches(0.85))
    grad_bar(slide, Inches(0.55), gy, Inches(5.5), Inches(0.055), color2, A_CYAN, angle=0)
    return gy + Inches(0.08)

def metric_tile(slide, x, y, value, label, color=A_CYAN, w=Inches(2.4), h=Inches(1.3)):
    s = rnd(slide, x, y, w, h, grad=(CARD_BG2, DARK_BG), adj=10000, shadow=True)
    # top color bar
    tb = rnd(slide, x + Inches(0.55), y + Inches(0.1), w - Inches(1.1), Inches(0.055),
             fill=color, adj=20000)
    txb(slide, value, x, y + Inches(0.2), w, Inches(0.7),
        size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, label, x, y + Inches(0.88), w, Inches(0.38),
        size=9, color=DIM, align=PP_ALIGN.CENTER)

def modern_card(slide, x, y, w, h, title, body_lines, icon="", accent=A_BLUE, adj=8000):
    s = rnd(slide, x, y, w, h, grad=(CARD_BG2, CARD_BG, 150), adj=adj, shadow=True)
    # left accent stripe
    rect(slide, x, y + Inches(0.18), Inches(0.048), h - Inches(0.36), fill=accent)
    tx = x + Inches(0.18)
    if icon:
        txb(slide, icon, tx, y + Inches(0.1), Inches(0.5), Inches(0.45),
            size=20, bold=True, color=accent)
        tx = x + Inches(0.65)
    txb(slide, title, tx, y + Inches(0.1), w - tx + x - Inches(0.18),
        Inches(0.42), size=12.5, bold=True, color=accent)
    by = y + Inches(0.58)
    tw = w - Inches(0.25) - (tx - x)
    for ln in body_lines:
        txb(slide, ln, tx, by, tw, Inches(0.32), size=10.5, color=LIGHT)
        by += Inches(0.31)

def arw_dn(slide, cx, y1, length=Inches(0.3), color=A_BLUE):
    connector(slide, cx, y1, cx, y1 + length, color, Pt(2))
    txb(slide, "▾", cx - Inches(0.13), y1 + length - Inches(0.07),
        Inches(0.26), Inches(0.26), size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

def arw_rt(slide, x1, cy, length=Inches(0.4), color=A_BLUE):
    connector(slide, x1, cy, x1 + length, cy, color, Pt(2))
    txb(slide, "▸", x1 + length - Inches(0.04), cy - Inches(0.13),
        Inches(0.26), Inches(0.26), size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0C, 0x16), RGBColor(0x0C, 0x1A, 0x2C), 145)
orbs(sl, A_BLUE)

# Left panel dark overlay
lp = rect(sl, 0, 0, Inches(7.1), H, fill=RGBColor(0x04, 0x0D, 0x1A))
lp.line.fill.background()

# Vertical grid lines decorative (right panel)
for i in range(7):
    rect(sl, Inches(7.2 + i * 0.87), 0, Pt(0.6), H, fill=RGBColor(0x15, 0x28, 0x3C))

# Horizontal grid lines (left panel)
for j in range(8):
    rect(sl, 0, Inches(j * 0.95), Inches(7.1), Pt(0.4), fill=RGBColor(0x0F, 0x1E, 0x2E))

# Top gradient bar
grad_bar(sl, 0, 0, Inches(7.1), Inches(0.07), A_BLUE, RGBColor(0x00, 0x55, 0x99), angle=0)

# Hackathon badge
pill(sl, "⬡  EPAM HACKATHON · STREAM 3", Inches(0.55), Inches(0.42),
     bg_c=RGBColor(0x00, 0x28, 0x44), fg_c=A_BLUE, size=9.5, w=Inches(3.6), h=Inches(0.3))

# Hero text
txb(sl, "Migration", Inches(0.55), Inches(1.0), Inches(6.2), Inches(1.35),
    size=78, bold=True, color=WHITE)
txb(sl, "Validator", Inches(0.55), Inches(2.2), Inches(6.2), Inches(1.35),
    size=78, bold=True, color=A_BLUE)

# Gradient underline
grad_bar(sl, Inches(0.55), Inches(3.5), Inches(5.2), Inches(0.065), A_BLUE, A_CYAN, angle=0)

# Subtitle
txb(sl, "Gemini-Enabled Enterprise Migration Intelligence",
    Inches(0.55), Inches(3.68), Inches(6.3), Inches(0.65),
    size=17.5, color=LIGHT)

# Quote
txb(sl, '"Natural language. Human oversight. Audit trail."',
    Inches(0.55), Inches(4.42), Inches(6.2), Inches(0.75),
    size=13, italic=True, color=DIM)

# Feature pills
for i, (t, c) in enumerate([("24 Gemini Tools", A_BLUE), ("RBAC · OCC · Audit", A_CYAN),
                              ("31/31 Tests Pass", A_GREEN)]):
    pill(sl, t, Inches(0.55 + i * 2.05), Inches(5.3), bg_c=RGBColor(0x0D, 0x22, 0x35),
         fg_c=c, size=9.5, w=Inches(1.9), h=Inches(0.3))

# Right — Architecture mini-diagram
txb(sl, "ARCHITECTURE", Inches(7.55), Inches(0.4), Inches(5.4), Inches(0.35),
    size=8.5, bold=True, color=A_BLUE)
connector(sl, Inches(7.55), Inches(0.7), Inches(12.8), Inches(0.7), A_BLUE, Pt(0.5))

arch = [
    ("⬡  Gemini", A_BLUE, RGBColor(0x00, 0x1E, 0x3A)),
    ("⧉  Migration Connector", A_CYAN, RGBColor(0x00, 0x1A, 0x2C)),
    ("⚙  Validation Engine", A_GREEN, RGBColor(0x00, 0x18, 0x18)),
    ("⛁  PostgreSQL · MSSQL · Athena", DIM, RGBColor(0x0E, 0x1C, 0x28)),
    ("❄  Snowflake Target", A_CYAN, RGBColor(0x0A, 0x1E, 0x30)),
]
ay = Inches(0.82)
for text, fg, fill_c in arch:
    s = rnd(sl, Inches(7.55), ay, Inches(5.38), Inches(0.65), fill=fill_c, adj=8000, shadow=True)
    s.fill.solid(); s.fill.fore_color.rgb = fill_c
    # left glow stripe
    gbar = rect(sl, Inches(7.55), ay + Inches(0.12), Inches(0.045), Inches(0.4), fill=fg)
    txb(sl, text, Inches(7.72), ay, Inches(5.0), Inches(0.65),
        size=13.5, bold=True, color=fg)
    if text != arch[-1][0]:
        arw_dn(sl, Inches(10.24), ay + Inches(0.65), Inches(0.22), A_BLUE)
    ay += Inches(0.9)

# Stats row
stats_s = [("24", "Gemini Tools"), ("5", "RBAC Roles"), ("3", "Source DBs"), ("31✓", "Security Tests")]
sx = Inches(7.55)
for val, lbl in stats_s:
    metric_tile(sl, sx, Inches(6.0), val, lbl, color=A_CYAN, w=Inches(1.32), h=Inches(1.12))
    sx += Inches(1.38)

footer(sl, "1 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x06, 0x0A, 0x14), RGBColor(0x0E, 0x18, 0x26), 140)
orbs(sl, A_RED)

section_tag(sl, "THE CHALLENGE", A_RED)
div_y = big_title(sl, "Enterprise migrations are", "slow, risky, and ungoverned.", A_RED)

# Left column — workflow steps
add_lv = rect(sl, Inches(0.45), div_y + Inches(0.22), Pt(2.5), Inches(4.8), fill=A_RED)

txb(sl, "TODAY'S WORKFLOW", Inches(0.75), div_y + Inches(0.22), Inches(5.5), Inches(0.35),
    size=9.5, bold=True, color=A_RED)

steps = [
    ("1", "Write column-mapping scripts manually", "Days of engineer time"),
    ("2", "Generate validation SQL per table",     "4–8 queries per table"),
    ("3", "Run validation queries",                "Manual execution per layer"),
    ("4", "Review mismatches in spreadsheets",     "No confidence signal"),
    ("5", "Escalate ambiguous cases to architects","Bottlenecks every time"),
    ("6", "No audit trail for decisions",          "Compliance exposure"),
]
sy = div_y + Inches(0.65)
for num, step, sub in steps:
    nb = rnd(sl, Inches(0.75), sy, Inches(0.44), Inches(0.44),
             fill=A_RED, adj=20000, shadow=True)
    nb.fill.solid(); nb.fill.fore_color.rgb = A_RED
    txb(sl, num, Inches(0.75), sy, Inches(0.44), Inches(0.44),
        size=13, bold=True, color=DARKER, align=PP_ALIGN.CENTER)
    txb(sl, step, Inches(1.32), sy + Inches(0.01), Inches(4.3), Inches(0.27),
        size=12, bold=True, color=WHITE)
    txb(sl, sub, Inches(1.32), sy + Inches(0.27), Inches(4.3), Inches(0.21),
        size=9.5, color=DIM)
    sy += Inches(0.68)

# Right column — cost cards
txb(sl, "THE COST AT SCALE", Inches(6.6), div_y + Inches(0.22), Inches(6.3), Inches(0.35),
    size=9.5, bold=True, color=DIM)

costs = [
    ("100–1000+", "Tables per migration project",           A_RED),
    ("4–8",       "Manual SQL queries per table per layer", A_ORANGE),
    ("0",         "Governed approval trail — none today",   A_RED),
    ("60–80%",    "Engineer time on mechanics, not design", A_ORANGE),
    ("Weeks",     "Validation time — should be hours",      DIM),
]
cy2 = div_y + Inches(0.65)
for val, desc, vc in costs:
    s = rnd(sl, Inches(6.6), cy2, Inches(6.28), Inches(0.9),
            grad=(RGBColor(0x18, 0x25, 0x38), CARD_BG), adj=8000, shadow=True)
    rect(sl, Inches(6.6), cy2 + Inches(0.14), Inches(0.045), Inches(0.62), fill=vc)
    txb(sl, val, Inches(6.78), cy2 + Inches(0.1), Inches(1.9), Inches(0.55),
        size=26, bold=True, color=vc)
    txb(sl, desc, Inches(8.75), cy2 + Inches(0.27), Inches(3.9), Inches(0.45),
        size=11, color=LIGHT)
    cy2 += Inches(1.0)

footer(sl, "2 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY GEMINI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0D, 0x18), RGBColor(0x0D, 0x1B, 0x2C), 135)
orbs(sl, A_BLUE)

section_tag(sl, "THE SOLUTION", A_BLUE)
div_y = big_title(sl, "Why Gemini changes everything",
                  "Natural language + governed tools.", A_BLUE)

# Comparison table
headers = ["CAPABILITY", "TRADITIONAL", "GEMINI-CONNECTED  ✓"]
col_x   = [Inches(0.45), Inches(4.2), Inches(8.6)]
col_w   = [Inches(3.7),  Inches(4.3), Inches(4.65)]
tbl_y   = div_y + Inches(0.25)

# Header row
for i, (hd, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg_c = A_BLUE if i == 2 else RGBColor(0x1E, 0x38, 0x50)
    fg_c = DARKER if i == 2 else WHITE
    s = rnd(sl, cx, tbl_y, cw, Inches(0.48), fill=bg_c, adj=6000)
    if i == 2:
        apply_grad(s, A_BLUE, RGBColor(0x00, 0x78, 0xB4), angle_deg=0)
        s.line.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = bg_c
    txb(sl, hd, cx, tbl_y, cw, Inches(0.48),
        size=10.5, bold=True, color=fg_c, align=PP_ALIGN.CENTER)

rows = [
    ("Querying data",       "Run scripts manually",             "Ask in natural language"),
    ("Interpreting results","Read raw SQL output",              "Get an explanation"),
    ("Ambiguous mappings",  "Engineer resolves (hours)",        "AI proposes → human approves"),
    ("Root-cause analysis", "Ad-hoc SQL debugging (hours)",     "Conversational (seconds)"),
    ("Governance",          "None — no trail, no versions",     "RBAC + OCC + audit trail"),
    ("Multi-source view",   "Separate scripts per system",      "Single federated question"),
]
ry = tbl_y + Inches(0.48)
for ri, (cap, trad, gem) in enumerate(rows):
    for ci, (cx, cw, txt) in enumerate(zip(col_x, col_w, [cap, trad, gem])):
        alt = RGBColor(0x11, 0x20, 0x30) if ri % 2 == 0 else CARD_BG
        bg_c = RGBColor(0x09, 0x1E, 0x30) if ci == 2 else alt
        s = rect(sl, cx, ry, cw, Inches(0.58), fill=bg_c)
        col = A_CYAN if ci == 2 else (WHITE if ci == 0 else DIM)
        txb(sl, ("✓  " if ci == 2 else "") + txt,
            cx + Inches(0.1), ry + Inches(0.14), cw - Inches(0.2), Inches(0.4),
            size=11.5, bold=(ci == 0), color=col)
    ry += Inches(0.58)

# Insight banner
banner_y = ry + Inches(0.2)
s = rnd(sl, Inches(0.45), banner_y, Inches(12.4), Inches(0.72),
        grad=(RGBColor(0x00, 0x2A, 0x44), RGBColor(0x00, 0x14, 0x28)), adj=8000)
rect(sl, Inches(0.45), banner_y, Inches(0.07), Inches(0.72), fill=A_BLUE)
txb(sl, "KEY INSIGHT", Inches(0.68), banner_y + Inches(0.1), Inches(1.5), Inches(0.28),
    size=9, bold=True, color=A_BLUE)
txb(sl, "Gemini drives 24 governed tools here — each permission-checked, version-controlled, and audit-logged. Not a chatbot.",
    Inches(2.3), banner_y + Inches(0.18), Inches(10.35), Inches(0.45),
    size=12.5, color=WHITE, italic=True)

footer(sl, "3 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0D, 0x18), RGBColor(0x0C, 0x1C, 0x2C), 130)
orbs(sl, A_CYAN)

section_tag(sl, "ARCHITECTURE", A_CYAN)
txb(sl, "Three-tier governed architecture",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8),
    size=38, bold=True, color=WHITE)
txb(sl, "Gemini  →  Connector  →  Engine  →  Databases",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.55),
    size=20, bold=False, color=A_CYAN)
grad_bar(sl, Inches(0.55), Inches(2.0), Inches(5.5), Inches(0.06), A_CYAN, A_BLUE, angle=0)

tiers = [
    ("TIER 1", "Gemini",
     "Natural language · Function calling · 10-round tool loop · Offline fallback",
     A_BLUE, RGBColor(0x00, 0x1A, 0x36)),
    ("TIER 2", "Migration Intelligence Connector  (FastAPI · Port 8001)",
     "Auth · RBAC (5 roles / 15 perms) · 24 Tools · OCC · Audit · Metrics",
     A_CYAN, RGBColor(0x00, 0x18, 0x26)),
    ("TIER 3", "Validation Engine",
     "CanonicalValidationPlan · SQL Generation · YAML Config · RuleBook",
     A_GREEN, RGBColor(0x00, 0x16, 0x14)),
]
ty = Inches(2.18)
for lbl, name, desc, accent, fill_c in tiers:
    s = rnd(sl, Inches(0.45), ty, Inches(7.8), Inches(1.12),
            fill=fill_c, adj=9000, shadow=True)
    s.fill.solid(); s.fill.fore_color.rgb = fill_c
    # Tier badge
    badge = rnd(sl, Inches(0.45), ty, Inches(1.15), Inches(1.12), fill=accent, adj=9000)
    badge.fill.solid(); badge.fill.fore_color.rgb = accent
    txb(sl, lbl, Inches(0.45), ty, Inches(1.15), Inches(1.12),
        size=9.5, bold=True, color=DARKER, align=PP_ALIGN.CENTER)
    txb(sl, name, Inches(1.76), ty + Inches(0.1), Inches(6.3), Inches(0.48),
        size=16, bold=True, color=WHITE)
    txb(sl, desc, Inches(1.76), ty + Inches(0.6), Inches(6.3), Inches(0.46),
        size=10.5, color=DIM)
    if lbl != "TIER 3":
        arw_dn(sl, Inches(4.35), ty + Inches(1.12), Inches(0.25), accent)
    ty += Inches(1.4)

# Source databases row
src_y = ty + Inches(0.26)
bg_src = rnd(sl, Inches(0.45), src_y - Inches(0.06), Inches(7.8), Inches(0.96),
             grad=(RGBColor(0x12, 0x20, 0x30), DARK_BG), adj=9000)
sources = [("PostgreSQL", A_BLUE), ("MSSQL", A_CYAN), ("AWS Athena", A_ORANGE), ("Snowflake", A_GREEN)]
sx2 = Inches(0.65)
for src, c in sources:
    s = rnd(sl, sx2, src_y, Inches(1.8), Inches(0.7), fill=DARK_BG, line=c, lw=Pt(1.2), adj=8000)
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BG
    txb(sl, src, sx2, src_y, Inches(1.8), Inches(0.7),
        size=12, bold=True, color=c, align=PP_ALIGN.CENTER)
    sx2 += Inches(1.96)

# Right callout cards
rx = Inches(8.55)
callouts = [
    ("24 TOOLS", ["Discovery · Mapping · Rules", "Plans · Execution · Write-back"], A_BLUE),
    ("5 ROLES · 15 PERMS",
     ["VIEWER  REVIEWER  RULE_ADMIN", "VALIDATION_OPERATOR  ADMIN"], A_CYAN),
    ("OCC VERSION CONTROL",
     ["check_and_bump() prevents", "concurrent write conflicts (HTTP 409)"], A_GREEN),
    ("APPEND-ONLY AUDIT",
     ["actor · timestamp · reason · version", "→ output/audit_log.jsonl"], A_ORANGE),
]
cy3 = Inches(2.18)
for title, body, accent in callouts:
    s = rnd(sl, rx, cy3, Inches(4.55), Inches(1.05),
            grad=(CARD_BG2, CARD_BG, 145), adj=9000, shadow=True)
    rect(sl, rx, cy3 + Inches(0.12), Inches(0.05), Inches(0.8), fill=accent)
    txb(sl, title, rx + Inches(0.2), cy3 + Inches(0.1), Inches(4.15), Inches(0.38),
        size=12, bold=True, color=accent)
    for i, ln in enumerate(body):
        txb(sl, ln, rx + Inches(0.2), cy3 + Inches(0.52) + i * Inches(0.28),
            Inches(4.15), Inches(0.28), size=10, color=LIGHT)
    cy3 += Inches(1.2)

footer(sl, "4 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LIVE USE CASE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0C, 0x14), RGBColor(0x0A, 0x18, 0x24), 130)
orbs(sl, A_GREEN)

section_tag(sl, "LIVE USE CASE", A_GREEN)
txb(sl, "From natural language to validated data",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.75), size=36, bold=True, color=WHITE)
txb(sl, "events table  ·  PostgreSQL (fms) → Snowflake Bronze",
    Inches(0.55), Inches(1.46), Inches(12.3), Inches(0.48), size=16, color=A_GREEN)
grad_bar(sl, Inches(0.55), Inches(1.9), Inches(5.5), Inches(0.055), A_GREEN, A_CYAN, angle=0)

# Chat transcript
conv_y = Inches(2.05)
conv = [
    ("USER",         '"Validate the events migration."',                                      A_CYAN,  False),
    ("TOOL CALL",    "→ discover_connections()",                                               A_BLUE,  True),
    ("RESULT",       "PostgreSQL (fms)  →  Snowflake (dev_edge_bronze)",                      DIM,     False),
    ("TOOL CALL",    "→ generate_validation_plan(table='events', layer='bronze')",             A_BLUE,  True),
    ("RESULT",       "38 columns mapped  ·  4 Fivetran cols excluded  ·  1 needs review (82%)", DIM,   False),
    ("GEMINI",       '"1 mapping needs review: created_ts → CREATED_AT  (82% confidence)"',   A_GREEN, False),
    ("USER",         '"I\'ve confirmed it in the DDL. Approve it."',                           A_CYAN,  False),
    ("TOOL CALL",    "→ approve_mapping(actor='jane@corp.com', reason='Confirmed DDL', expected_version=0)", A_BLUE, True),
    ("RESULT",       "APPROVED  ·  Version = 1  ·  AuditRecord written",                      A_GREEN, False),
    ("TOOL CALL",    "→ execute_validation(actor='jane@corp.com')",                            A_BLUE,  True),
    ("RESULT",       "STATUS: PASS  ·  Coverage: 100%  ·  38 columns  ·  0 failures",         A_GREEN, False),
    ("GEMINI",       '"All 38 columns validated. No mismatches detected."',                    A_GREEN, False),
]
for role, text, color, is_tool in conv:
    label_bg = RGBColor(0x00, 0x25, 0x44) if is_tool else RGBColor(0x12, 0x22, 0x32)
    row_bg   = RGBColor(0x00, 0x12, 0x22) if is_tool else RGBColor(0x0E, 0x1C, 0x2C)
    # Role label
    s = rnd(sl, Inches(0.45), conv_y, Inches(1.38), Inches(0.33),
            fill=label_bg, adj=6000)
    s.fill.solid(); s.fill.fore_color.rgb = label_bg
    txb(sl, role, Inches(0.45), conv_y, Inches(1.38), Inches(0.33),
        size=8.5, bold=True, color=A_BLUE if is_tool else DIM, align=PP_ALIGN.CENTER)
    # Text row
    s2 = rnd(sl, Inches(1.88), conv_y, Inches(11.0), Inches(0.33),
             fill=row_bg, adj=6000)
    s2.fill.solid(); s2.fill.fore_color.rgb = row_bg
    txb(sl, text, Inches(2.0), conv_y, Inches(10.75), Inches(0.33),
        size=10.5, color=color, bold=is_tool)
    conv_y += Inches(0.37)

# Bottom stat tiles
stat_y = H - Inches(1.55)
grad_bar(sl, Inches(0.45), stat_y - Inches(0.12), Inches(12.45), Inches(0.05), A_BLUE, A_GREEN, angle=0)
stats3 = [("4", "Tool calls by Gemini"), ("0", "SQL written by hand"),
          ("1", "Human decision logged"), ("< 30s", "Total time"), ("100%", "Coverage")]
sx3 = Inches(0.45)
for val, lbl in stats3:
    metric_tile(sl, sx3, stat_y, val, lbl, color=A_GREEN, w=Inches(2.42), h=Inches(1.22))
    sx3 += Inches(2.57)

footer(sl, "5 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HUMAN GOVERNANCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x06, 0x0C, 0x12), RGBColor(0x0C, 0x18, 0x26), 130)
orbs(sl, A_ORANGE)

section_tag(sl, "HUMAN-IN-THE-LOOP", A_ORANGE)
txb(sl, "AI proposes. Humans decide.",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8), size=38, bold=True, color=WHITE)
txb(sl, "Every decision is on record.",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.55), size=22, bold=True, color=A_ORANGE)
grad_bar(sl, Inches(0.55), Inches(2.0), Inches(5.5), Inches(0.055), A_ORANGE, A_CYAN, angle=0)

# Confidence tiers (left)
txb(sl, "CONFIDENCE ROUTING", Inches(0.55), Inches(2.18), Inches(5.6), Inches(0.34),
    size=9.5, bold=True, color=A_ORANGE)

conf_rows = [
    ("≥ 95%",  "AUTO-ACCEPTED",    "No human review needed",         A_GREEN),
    ("75–94%", "STANDARD REVIEW",  "Queued in ApprovalStore",        A_ORANGE),
    ("< 75%",  "MANDATORY REVIEW", "Never auto-accepted — escalated", A_RED),
]
cy4 = Inches(2.6)
for threshold, action, note, color in conf_rows:
    # Threshold pill
    tp = rnd(sl, Inches(0.55), cy4, Inches(1.35), Inches(0.95),
             fill=color, adj=12000, shadow=True)
    tp.fill.solid(); tp.fill.fore_color.rgb = color
    txb(sl, threshold, Inches(0.55), cy4, Inches(1.35), Inches(0.95),
        size=19, bold=True, color=DARKER, align=PP_ALIGN.CENTER)
    # Action + note
    s = rnd(sl, Inches(2.02), cy4, Inches(4.1), Inches(0.95),
            grad=(CARD_BG2, CARD_BG, 150), adj=10000, shadow=True)
    rect(sl, Inches(2.02), cy4 + Inches(0.12), Inches(0.04), Inches(0.7), fill=color)
    txb(sl, action, Inches(2.16), cy4 + Inches(0.1), Inches(3.8), Inches(0.38),
        size=14, bold=True, color=color)
    txb(sl, note, Inches(2.16), cy4 + Inches(0.52), Inches(3.8), Inches(0.36),
        size=10.5, color=DIM)
    cy4 += Inches(1.08)

# AI self-approval guard
guard_y = cy4 + Inches(0.08)
gs = rnd(sl, Inches(0.55), guard_y, Inches(5.57), Inches(0.82),
         grad=(RGBColor(0x30, 0x08, 0x08), RGBColor(0x20, 0x05, 0x05)), adj=9000)
rect(sl, Inches(0.55), guard_y + Inches(0.1), Inches(0.045), Inches(0.62), fill=A_RED)
txb(sl, "⚠  AI SELF-APPROVAL BLOCKED",
    Inches(0.72), guard_y + Inches(0.1), Inches(5.0), Inches(0.35),
    size=12, bold=True, color=A_RED)
txb(sl, '"gemini_ai" and "ai" actor strings rejected at the tool level.',
    Inches(0.72), guard_y + Inches(0.44), Inches(5.0), Inches(0.3),
    size=10.5, color=LIGHT)

# Approval flow (right)
txb(sl, "APPROVAL FLOW", Inches(6.6), Inches(2.18), Inches(6.3), Inches(0.34),
    size=9.5, bold=True, color=A_ORANGE)

flow = [
    ("Gemini Recommendation",          A_BLUE,   RGBColor(0x00, 0x1E, 0x38)),
    ("ApprovalStore  {status: PENDING}", A_ORANGE, RGBColor(0x30, 0x28, 0x08)),
    ("Human Reviews Decision",          A_ORANGE,  RGBColor(0x28, 0x18, 0x04)),
]
fy = Inches(2.6)
fx = Inches(7.3)
fw = Inches(4.6)
for text, accent, fill_c in flow:
    s = rnd(sl, fx, fy, fw, Inches(0.56), fill=fill_c, adj=9000, shadow=True)
    s.fill.solid(); s.fill.fore_color.rgb = fill_c
    rect(sl, fx, fy + Inches(0.08), Inches(0.04), Inches(0.4), fill=accent)
    txb(sl, text, fx + Inches(0.14), fy, fw - Inches(0.14), Inches(0.56),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    arw_dn(sl, fx + fw/2, fy + Inches(0.56), Inches(0.25), A_ORANGE)
    fy += Inches(0.84)

# Outcome branches
outcomes = [("Approve", A_GREEN), ("Modify", A_ORANGE), ("Reject", A_RED)]
ox = Inches(6.68)
for label, color in outcomes:
    s = rnd(sl, ox, fy, Inches(1.85), Inches(0.54), fill=color, adj=10000, shadow=True)
    s.fill.solid(); s.fill.fore_color.rgb = color
    txb(sl, label, ox, fy, Inches(1.85), Inches(0.54),
        size=13, bold=True, color=DARKER, align=PP_ALIGN.CENTER)
    ox += Inches(2.08)
fy += Inches(0.6)
arw_dn(sl, fx + fw/2, fy, Inches(0.22), A_GREEN)
fy += Inches(0.25)
fbox = rnd(sl, fx, fy, fw, Inches(0.52), fill=DARK_BG, line=A_GREEN, lw=Pt(1.2), adj=9000)
fbox.fill.solid(); fbox.fill.fore_color.rgb = DARK_BG
txb(sl, "CanonicalValidationPlan  →  SQL  →  Execution",
    fx, fy, fw, Inches(0.52), size=11.5, color=A_GREEN, align=PP_ALIGN.CENTER)
fy += Inches(0.62)
txb(sl, "actor · timestamp · reason · version → audit_log.jsonl",
    Inches(6.7), fy, Inches(6.0), Inches(0.35), size=10.5, italic=True, color=DIM)

footer(sl, "6 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x06, 0x08, 0x14), RGBColor(0x0D, 0x16, 0x26), 125)
orbs(sl, A_RED)

section_tag(sl, "SECURITY", A_RED)
txb(sl, "Enterprise-grade security,",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8), size=38, bold=True, color=WHITE)
txb(sl, "built in from day one.",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.8), size=38, bold=True, color=A_RED)
grad_bar(sl, Inches(0.55), Inches(2.2), Inches(5.5), Inches(0.055), A_RED, A_ORANGE, angle=0)

layers = [
    ("1", "AUTHENTICATION",
     ["JWT (HS256) — enterprise SSO",
      "Static token — CI/CD pipelines",
      "Dev mode — local only (ADMIN, no validation)"],
     A_BLUE, "🔐"),
    ("2", "AUTHORIZATION (RBAC)",
     ["5 roles: VIEWER → REVIEWER → RULE_ADMIN",
      "           VALIDATION_OPERATOR → ADMIN",
      "15 fine-grained permissions",
      "Resource-level allowlists (source / DB / schema)"],
     A_CYAN, "🛡"),
    ("3", "CONCURRENCY CONTROL",
     ["Optimistic concurrency (VersionStore)",
      "check_and_bump(entity_key, expected_version)",
      "HTTP 409 VERSION_CONFLICT on stale write",
      "Prevents parallel approval conflicts"],
     A_ORANGE, "🔄"),
    ("4", "AUDIT TRAIL",
     ["Append-only JSONL (never modified or deleted)",
      "actor · user_id · timestamp · reason · version",
      "No secrets, passwords, or raw data logged",
      "request_id links to tool call log"],
     A_GREEN, "📋"),
]

card_y = Inches(2.38)
lx = Inches(0.45)
card_w = Inches(3.15)
for num, title, bullets, accent, icon in layers:
    # Number badge
    nb = rnd(sl, lx, card_y, Inches(0.46), Inches(0.46), fill=accent, adj=20000, shadow=True)
    nb.fill.solid(); nb.fill.fore_color.rgb = accent
    txb(sl, num, lx, card_y, Inches(0.46), Inches(0.46),
        size=14, bold=True, color=DARKER, align=PP_ALIGN.CENTER)
    # Card
    s = rnd(sl, lx, card_y + Inches(0.44), card_w, Inches(4.55),
            grad=(CARD_BG2, CARD_BG, 150), adj=10000, shadow=True)
    rect(sl, lx, card_y + Inches(0.44), card_w, Inches(0.055), fill=accent)
    # Icon
    txb(sl, icon, lx + Inches(0.15), card_y + Inches(0.55), Inches(0.5), Inches(0.48),
        size=22, bold=True, color=accent)
    txb(sl, title, lx + Inches(0.7), card_y + Inches(0.58), card_w - Inches(0.85), Inches(0.42),
        size=12, bold=True, color=accent)
    by2 = card_y + Inches(1.1)
    for b in bullets:
        txb(sl, "▸  " + b, lx + Inches(0.18), by2, card_w - Inches(0.36), Inches(0.34),
            size=10.5, color=LIGHT)
        by2 += Inches(0.36)
    lx += Inches(3.3)

# Zero-credential banner
bny = card_y + Inches(5.1)
bs = rnd(sl, Inches(0.45), bny, Inches(12.43), Inches(0.75),
         grad=(RGBColor(0x06, 0x1E, 0x06), RGBColor(0x04, 0x14, 0x04)), adj=9000)
rect(sl, Inches(0.45), bny, Inches(0.06), Inches(0.75), fill=A_GREEN)
txb(sl, "ZERO CREDENTIAL EXPOSURE",
    Inches(0.68), bny + Inches(0.1), Inches(2.8), Inches(0.28),
    size=10, bold=True, color=A_GREEN)
txb(sl, "Gemini receives tool results only — never passwords, connection strings, or API keys.   "
        "31 / 31 security tests pass.",
    Inches(3.65), bny + Inches(0.2), Inches(9.0), Inches(0.45),
    size=12.5, color=WHITE)

footer(sl, "7 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RULE BOOK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0D, 0x1A), RGBColor(0x0C, 0x1C, 0x2E), 130)
orbs(sl, A_CYAN)

section_tag(sl, "RULE BOOK", A_CYAN)
txb(sl, "Deterministic validation,",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8), size=38, bold=True, color=WHITE)
txb(sl, "AI-augmented governance.",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.8), size=38, bold=True, color=A_CYAN)
grad_bar(sl, Inches(0.55), Inches(2.2), Inches(5.5), Inches(0.055), A_CYAN, A_BLUE, angle=0)

# Rule category cards (left)
cats = [
    ("Static Exclusions",    "Built-in Fivetran audit columns excluded automatically — no config", A_RED,    "BUILT-IN",    "🔒"),
    ("Pattern Rules",        "Regex: ^_FIVETRAN_.* matches entire column families",               A_ORANGE,  "CONFIG",       "🔍"),
    ("Global Exclusions",    "User-defined per DB type — persisted to YAML for all tables",       A_BLUE,    "USER-DEFINED", "⚙"),
    ("Transformation Rules", "boolean→'0'/'1', timestamp→text, hstore→json, uuid→text",          A_CYAN,    "GENERATED",    "⚡"),
    ("Learned Rules",        "AI-proposed, human-reviewed (RULE_ADMIN only), versioned",          A_GREEN,   "AI + HUMAN",   "🤖"),
]
cry = Inches(2.38)
for name, desc, color, badge, icon in cats:
    s = rnd(sl, Inches(0.45), cry, Inches(6.25), Inches(0.88),
            grad=(CARD_BG2, CARD_BG, 150), adj=9000, shadow=True)
    rect(sl, Inches(0.45), cry, Inches(0.048), Inches(0.88), fill=color)
    # Icon circle background
    ic = rnd(sl, Inches(0.6), cry + Inches(0.2), Inches(0.48), Inches(0.48),
             fill=DARKER, adj=20000)
    ic.fill.solid(); ic.fill.fore_color.rgb = DARKER
    txb(sl, icon, Inches(0.6), cry + Inches(0.2), Inches(0.48), Inches(0.48),
        size=18, align=PP_ALIGN.CENTER)
    pill(sl, badge, Inches(5.48), cry + Inches(0.32), bg_c=color, fg_c=DARKER,
         size=8.5, w=Inches(1.1), h=Inches(0.24))
    txb(sl, name, Inches(1.2), cry + Inches(0.08), Inches(4.1), Inches(0.36),
        size=13, bold=True, color=color)
    txb(sl, desc, Inches(1.2), cry + Inches(0.48), Inches(4.2), Inches(0.34),
        size=10, color=DIM)
    cry += Inches(0.97)

# Lifecycle (right)
txb(sl, "LEARNED RULE LIFECYCLE", Inches(7.05), Inches(2.38), Inches(5.8), Inches(0.34),
    size=9.5, bold=True, color=A_CYAN)

lifecycle = [
    ("AI Proposal",  "Observed pattern during pipeline run",  A_BLUE),
    ("Draft",        "Saved to rule_book_learned.json",        RGBColor(0x88, 0x88, 0x10)),
    ("Human Review", "RULE_ADMIN reviews and decides",          A_ORANGE),
    ("Approved",     "Queued for activation",                  A_CYAN),
    ("Active",       "Applied to all future pipeline runs",    A_GREEN),
    ("Versioned",    "Every change tracked in VersionStore",   DIM),
]
lfy = Inches(2.82)
for lname, ldesc, lcolor in lifecycle:
    ls = rnd(sl, Inches(7.05), lfy, Inches(5.82), Inches(0.57),
             fill=DARK_BG, line=lcolor, lw=Pt(1.2), adj=9000, shadow=True)
    ls.fill.solid(); ls.fill.fore_color.rgb = DARK_BG
    rect(sl, Inches(7.05), lfy + Inches(0.09), Inches(0.045), Inches(0.38), fill=lcolor)
    txb(sl, lname, Inches(7.22), lfy + Inches(0.08), Inches(2.1), Inches(0.3),
        size=12.5, bold=True, color=lcolor)
    txb(sl, ldesc, Inches(9.4), lfy + Inches(0.14), Inches(3.35), Inches(0.36),
        size=10, color=DIM)
    if lname != "Versioned":
        arw_dn(sl, Inches(9.96), lfy + Inches(0.57), Inches(0.2), lcolor)
    lfy += Inches(0.78)

footer(sl, "8 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0C, 0x14), RGBColor(0x0A, 0x18, 0x24), 125)
orbs(sl, A_GREEN)

section_tag(sl, "BUSINESS IMPACT", A_GREEN)
txb(sl, "Measurable value from day one.",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8), size=38, bold=True, color=WHITE)
txb(sl, "Prototype benchmarks — controlled demo environment",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.45), size=14, color=DIM, italic=True)
grad_bar(sl, Inches(0.55), Inches(1.9), Inches(5.5), Inches(0.055), A_GREEN, A_CYAN, angle=0)

# Big metric tiles
tiles = [
    ("100%",  "Columns auto-mapped\n(events table)",    A_GREEN),
    ("4",     "Manual SQL queries\navoided per table",  A_CYAN),
    ("< 30s", "Plan generation\ntime per table",         A_BLUE),
    ("31/31", "Security tests\npassing",                 A_GREEN),
    ("0",     "Credential leaks\ndetected",              A_GREEN),
]
tx2 = Inches(0.45)
for val, lbl, vc in tiles:
    metric_tile(sl, tx2, Inches(2.08), val, lbl, color=vc, w=Inches(2.48), h=Inches(1.3))
    tx2 += Inches(2.63)

# Scale projection (left)
ps = rnd(sl, Inches(0.45), Inches(3.55), Inches(5.9), Inches(3.5),
         grad=(CARD_BG2, DARK_BG, 150), adj=10000, shadow=True)
rect(sl, Inches(0.45), Inches(3.55), Inches(5.9), Inches(0.055), fill=A_CYAN)
txb(sl, "SCALE PROJECTION  (illustrative)", Inches(0.65), Inches(3.62),
    Inches(5.5), Inches(0.32), size=9.5, bold=True, color=A_CYAN)

proj = [
    ("200 tables migrated",           "~6,300 engineer hours saved"),
    ("40 columns avg · 90% auto",     "~90% reduction in validation effort"),
    ("3 source systems",              "Single Gemini interface — fully federated"),
    ("€100 / hr blended rate",        "~€630,000 saved per migration project"),
]
py2 = Inches(4.05)
for left, right in proj:
    txb(sl, "▸  " + left, Inches(0.65), py2, Inches(2.75), Inches(0.38), size=11, color=LIGHT)
    txb(sl, right, Inches(3.5), py2, Inches(2.6), Inches(0.38), size=11, bold=True, color=A_CYAN)
    py2 += Inches(0.46)

txb(sl, "* Illustrative — based on industry-typical effort estimates",
    Inches(0.65), py2 + Inches(0.1), Inches(5.5), Inches(0.28), size=9, italic=True, color=DIM)

# Value pillars (right)
ps2 = rnd(sl, Inches(6.6), Inches(3.55), Inches(6.3), Inches(3.5),
          grad=(CARD_BG2, DARK_BG, 145), adj=10000, shadow=True)
rect(sl, Inches(6.6), Inches(3.55), Inches(6.3), Inches(0.055), fill=A_GREEN)
txb(sl, "ENTERPRISE VALUE PROPOSITIONS", Inches(6.8), Inches(3.62),
    Inches(5.9), Inches(0.32), size=9.5, bold=True, color=A_GREEN)

pillars = [
    ("Compliance", "Every mapping: auditable human actor, timestamp, reason, version.", A_GREEN),
    ("Risk",       "Confidence scoring surfaces ambiguous mappings before production.",  A_ORANGE),
    ("Speed",      "Gemini explains failures in natural language — no SQL expertise.",   A_BLUE),
    ("Scale",      "Portfolio view across all layers in one natural language question.", A_CYAN),
]
vpy = Inches(4.05)
for pt, desc, pc in pillars:
    rect(sl, Inches(6.75), vpy, Inches(0.045), Inches(0.72), fill=pc)
    txb(sl, pt, Inches(6.95), vpy + Inches(0.04), Inches(1.3), Inches(0.28),
        size=13, bold=True, color=pc)
    txb(sl, desc, Inches(8.35), vpy + Inches(0.1), Inches(4.35), Inches(0.55),
        size=10.5, color=LIGHT)
    vpy += Inches(0.84)

footer(sl, "9 / 10")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ROADMAP & CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, RGBColor(0x05, 0x0D, 0x1C), RGBColor(0x0C, 0x1C, 0x30), 140)
orbs(sl, A_BLUE)

section_tag(sl, "ROADMAP", A_BLUE)
txb(sl, "Built for scale. Ready to extend.",
    Inches(0.55), Inches(0.75), Inches(12.3), Inches(0.8), size=38, bold=True, color=WHITE)
txb(sl, "The connector pattern is extensible — any tool, any source, any governance.",
    Inches(0.55), Inches(1.5), Inches(12.3), Inches(0.48), size=15, color=LIGHT)
grad_bar(sl, Inches(0.55), Inches(1.95), Inches(5.5), Inches(0.055), A_BLUE, A_CYAN, angle=0)

# Current state card
cs = rnd(sl, Inches(0.45), Inches(2.12), Inches(6.1), Inches(4.85),
         grad=(CARD_BG2, DARK_BG, 145), adj=10000, shadow=True)
rect(sl, Inches(0.45), Inches(2.12), Inches(6.1), Inches(0.055), fill=A_GREEN)
pill(sl, "✓  CURRENT  ·  PROTOTYPE", Inches(0.65), Inches(2.18),
     bg_c=RGBColor(0x05, 0x20, 0x0D), fg_c=A_GREEN, size=9, w=Inches(2.5), h=Inches(0.28))

current = [
    "PostgreSQL · MSSQL · AWS Athena  →  Snowflake",
    "24 Gemini function-calling tools (6 categories)",
    "5-role RBAC + JWT / Static token authentication",
    "Optimistic concurrency (VersionStore · HTTP 409)",
    "Append-only audit log (AuditLogger · JSONL)",
    "Confidence-based human review routing",
    "Streamlit Web UI + FastAPI Connector + CLI",
    "31 / 31 security tests passing",
]
cy5 = Inches(2.6)
for item in current:
    txb(sl, "✓  " + item, Inches(0.72), cy5, Inches(5.5), Inches(0.38),
        size=11.5, color=LIGHT)
    cy5 += Inches(0.43)

# Planned card
ps3 = rnd(sl, Inches(6.82), Inches(2.12), Inches(6.1), Inches(4.85),
          grad=(CARD_BG2, DARK_BG, 145), adj=10000, shadow=True)
rect(sl, Inches(6.82), Inches(2.12), Inches(6.1), Inches(0.055), fill=A_ORANGE)
pill(sl, "◦  PLANNED", Inches(7.02), Inches(2.18),
     bg_c=RGBColor(0x28, 0x14, 0x04), fg_c=A_ORANGE, size=9, w=Inches(1.4), h=Inches(0.28))

planned = [
    ("Production hardening",        "Connection pooling · Async execution"),
    ("Additional sources",          "Oracle · BigQuery · MySQL · DB2"),
    ("Streaming validation",        "Real-time CDC change detection"),
    ("Enterprise SSO",              "OIDC / OAuth 2.0 integration"),
    ("Slack notifications",         "Approval requests routed to reviewers"),
    ("Multi-tenant deployment",     "Per-project namespace isolation"),
    ("AI remediation",              "Auto-fix suggestions from Gemini"),
]
py3 = Inches(2.6)
for title, desc in planned:
    txb(sl, "◦  " + title, Inches(7.1), py3, Inches(2.75), Inches(0.28),
        size=12, bold=True, color=A_ORANGE)
    txb(sl, desc, Inches(7.1), py3 + Inches(0.27), Inches(5.6), Inches(0.22),
        size=10, color=DIM)
    py3 += Inches(0.64)

# Closing tagline
tag_y = Inches(7.05)
ct = rnd(sl, Inches(0.45), tag_y - Inches(0.38), Inches(12.43), Inches(0.78),
         grad=(RGBColor(0x00, 0x18, 0x38), DARK_BG, 0), adj=9000)
rect(sl, Inches(0.45), tag_y - Inches(0.38), Inches(0.07), Inches(0.78), fill=A_BLUE)
txb(sl, "Migration Validator",
    Inches(0.72), tag_y - Inches(0.32), Inches(3.0), Inches(0.3),
    size=15, bold=True, color=A_BLUE)
txb(sl, "Governed enterprise migration validation through natural language.  "
        "Human oversight. Audit trail. Production-ready.",
    Inches(3.85), tag_y - Inches(0.28), Inches(9.0), Inches(0.5),
    size=13, italic=True, color=WHITE)

footer(sl, "10 / 10")


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Migration_Validator_Deck.pptx")
prs.save(out)
print(f"Saved  →  {out}")
print(f"Slides :  {len(prs.slides)}")
